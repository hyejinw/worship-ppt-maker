import io
import httpx
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

SLIDE_WIDTH = Inches(13.33)   # 33.87cm
SLIDE_HEIGHT = Inches(7.5)    # 19.05cm

FONT_SPECS = {
    "NanumGothic": {"file": "NanumGothic.ttf", "family": "NanumGothic"},
    "NanumMyeongjo": {"file": "NanumMyeongjo.ttf", "family": "NanumMyeongjo"},
    "NanumSquare": {"file": "NanumSquareR.ttf", "family": "NanumSquare"},
    "NotoSansKR": {"file": "NotoSansKR-Regular.ttf", "family": "Noto Sans KR"},
    "ATitleGothic1": {"file": "a타이틀고딕1.ttf", "family": "a타이틀고딕1"},
    "ATitleGothic2": {"file": "a타이틀고딕2.ttf", "family": "a타이틀고딕2"},
    "ATitleGothic3": {"file": "a타이틀고딕3.ttf", "family": "a타이틀고딕3"},
}


def _get_font_spec(font_family: str) -> dict[str, str]:
    return FONT_SPECS.get(font_family, FONT_SPECS["NanumGothic"])


def _apply_font_family(font, font_family: str):
    family_name = _get_font_spec(font_family)["family"]
    font.name = family_name

    r_pr = font._rPr
    for tag in ("latin", "ea", "cs"):
        node = r_pr.find(qn(f"a:{tag}"))
        if node is None:
            node = etree.SubElement(r_pr, qn(f"a:{tag}"))
        node.set("typeface", family_name)


def _hex_to_rgb(hex_color: str) -> tuple[int, int, int]:
    hex_color = hex_color.lstrip("#")
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))


def _add_solid_background(slide, hex_color: str):
    r, g, b = _hex_to_rgb(hex_color)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)


def _add_image_background(slide, image_data: bytes):
    pic_path = io.BytesIO(image_data)
    slide.shapes.add_picture(
        pic_path,
        left=0,
        top=0,
        width=SLIDE_WIDTH,
        height=SLIDE_HEIGHT,
    )


def _add_overlay(slide, opacity: float):
    if opacity <= 0:
        return
    shape = slide.shapes.add_shape(
        1,
        0, 0, SLIDE_WIDTH, SLIDE_HEIGHT,
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)

    alpha = int((1 - opacity) * 100000)
    sp_tree = shape._element
    sp_pr = sp_tree.find(qn("p:spPr"))
    solid_fill = sp_pr.find(qn("a:solidFill"))
    if solid_fill is not None:
        srgb_clr = solid_fill.find(qn("a:srgbClr"))
        if srgb_clr is None:
            srgb_clr = etree.SubElement(solid_fill, qn("a:srgbClr"))
            srgb_clr.set("val", "000000")
        alpha_elem = etree.SubElement(srgb_clr, qn("a:alpha"))
        alpha_elem.set("val", str(100000 - alpha))

    line = shape.line
    line.fill.background()


def _add_text_box(
    slide,
    lyrics: str,
    font_family: str,
    font_size: int,
    x_pct: float,
    y_pct: float,
    font_color: str = "#ffffff",
    text_box_width_pct: float = 95,
):
    line_count = len([l for l in lyrics.split("\n") if l.strip()])
    width_ratio = max(0.35, min(0.95, text_box_width_pct / 100))
    box_width = Emu(int(SLIDE_WIDTH * width_ratio))
    line_height_pt = font_size * 1.5
    box_height = Pt(line_height_pt * max(line_count, 1) + font_size)

    left = (SLIDE_WIDTH - box_width) / 2
    top_center = SLIDE_HEIGHT * (y_pct / 100) - box_height / 2
    top = max(0, min(top_center, SLIDE_HEIGHT - box_height))

    txBox = slide.shapes.add_textbox(left, top, box_width, box_height)
    tf = txBox.text_frame
    tf.word_wrap = True

    r, g, b = _hex_to_rgb(font_color)

    lines = lyrics.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = line

        font = run.font
        font.size = Pt(font_size)
        font.color.rgb = RGBColor(r, g, b)
        font.bold = False

        _apply_font_family(font, font_family)


def _add_title_box(
    slide,
    title: str,
    font_family: str,
    font_size: int,
    x_pct: float,
    y_pct: float,
    font_color: str = "#ffffff",
):
    title_font_size = max(12, int(round(font_size * 0.48)))
    box_width = Emu(int(SLIDE_WIDTH * 0.32))
    box_height = Pt(title_font_size * 1.8)

    left_center = SLIDE_WIDTH * (x_pct / 100) - box_width / 2
    top_center = SLIDE_HEIGHT * (y_pct / 100) - box_height / 2
    left = max(0, min(left_center, SLIDE_WIDTH - box_width))
    top = max(0, min(top_center, SLIDE_HEIGHT - box_height))

    txBox = slide.shapes.add_textbox(left, top, box_width, box_height)
    tf = txBox.text_frame
    tf.word_wrap = False

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"- {title.strip()}"

    font = run.font
    r, g, b = _hex_to_rgb(font_color)
    font.size = Pt(title_font_size)
    font.color.rgb = RGBColor(r, g, b)
    font.bold = False

    _apply_font_family(font, font_family)


def _apply_background(slide, bg_type, bg_value, bg_image_data):
    if bg_type == "black":
        _add_solid_background(slide, "#000000")
    elif bg_type == "color" and bg_value:
        _add_solid_background(slide, bg_value)
    elif bg_type == "image" and bg_image_data:
        _add_image_background(slide, bg_image_data)
    else:
        _add_solid_background(slide, "#000000")


async def build_pptx(
    slides: list[dict],
    settings: dict,
    songs_data: list[dict] | None = None,
    merge_songs: bool = True,
    export_song_id: str | None = None,
) -> bytes:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    blank_layout = prs.slide_layouts[6]

    font_family = settings.get("font_family", "NanumGothic")
    font_size = settings.get("font_size", 40)
    font_color = settings.get("font_color", "#ffffff")
    text_pos = settings.get("text_position", {"x": 50, "y": 30})
    title_pos = settings.get("title_position", {"x": 86, "y": 92})
    text_box_width = settings.get("text_box_width", 95)
    x_pct = text_pos.get("x", 50)
    y_pct = text_pos.get("y", 30)
    title_x_pct = title_pos.get("x", 86)
    title_y_pct = title_pos.get("y", 92)
    bg_type = settings.get("bg_type", "black")
    bg_value = settings.get("bg_value")
    overlay_opacity = settings.get("overlay_opacity", 0.0)
    show_title = settings.get("show_title", True)

    bg_image_data: bytes | None = None
    if bg_type == "image" and bg_value and bg_value.startswith("http"):
        async with httpx.AsyncClient(timeout=10) as client:
            resp = await client.get(bg_value)
            if resp.status_code == 200:
                bg_image_data = resp.content

    song_title_map: dict[str, str] = {}
    if songs_data:
        for s in songs_data:
            song_title_map[s["id"]] = s["title"]

    sorted_slides = sorted(slides, key=lambda s: s["order"])

    def add_slide(lyrics: str, title: str | None = None):
        slide = prs.slides.add_slide(blank_layout)
        _apply_background(slide, bg_type, bg_value, bg_image_data)
        if overlay_opacity > 0:
            _add_overlay(slide, overlay_opacity)
        if lyrics.strip():
            _add_text_box(slide, lyrics, font_family, font_size, x_pct, y_pct, font_color, text_box_width)
        if show_title and title:
            _add_title_box(slide, title, font_family, font_size, title_x_pct, title_y_pct, font_color)

    separator_slides = settings.get("separator_slides", True)

    if merge_songs and songs_data:
        song_ids = [s["id"] for s in songs_data]
        grouped: dict[str, list[dict]] = {sid: [] for sid in song_ids}
        ungrouped: list[dict] = []

        for s in sorted_slides:
            sid = s.get("song_id")
            if sid and sid in grouped:
                grouped[sid].append(s)
            else:
                ungrouped.append(s)

        if ungrouped:
            # song_id 없는 슬라이드: 순서대로 출력
            for slide_data in sorted_slides:
                lyrics_text = slide_data.get("lyrics", "").strip()
                add_slide(lyrics_text)
        else:
            for i, sid in enumerate(song_ids):
                song_slides = grouped.get(sid, [])
                title = song_title_map.get(sid)
                for slide_data in song_slides:
                    lyrics_text = slide_data.get("lyrics", "").strip()
                    add_slide(lyrics_text, title)
                # 곡 사이 빈 슬라이드 (마지막 곡 제외)
                if separator_slides and i < len(song_ids) - 1:
                    add_slide("")
    else:
        # 개별 곡 export
        if export_song_id:
            target_slides = [s for s in sorted_slides if s.get("song_id") == export_song_id]
            target_title = song_title_map.get(export_song_id) if show_title else None
        else:
            target_slides = sorted_slides
            target_title = None

        for slide_data in target_slides:
            lyrics_text = slide_data.get("lyrics", "").strip()
            add_slide(lyrics_text, target_title)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
